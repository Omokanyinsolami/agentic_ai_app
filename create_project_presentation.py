from __future__ import annotations

import json
import os
from functools import lru_cache
from pathlib import Path
from typing import Iterable

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
RESULTS_PATH = ROOT / "Backend" / "evaluation_results.json"
OUTPUT_PATH = ROOT / "Agentic_AI_for_Education_Presentation.pptx"
NOTES_PATH = ROOT / "Agentic_AI_for_Education_Presentation_Notes.md"
PROTOTYPE_OUTPUT_PATH = ROOT / "prototype" / OUTPUT_PATH.name
ASSETS_DIR = ROOT / "presentation_assets"

ARCH_PATH = ROOT / "Latex" / "prototype_architecture.png"
USE_CASE_PATH = ROOT / "Latex" / "dissertation_figures" / "figure_use_case_overview.png"
REQUEST_PATH = ROOT / "Latex" / "dissertation_figures" / "figure_request_workflow.png"
REMINDER_PATH = ROOT / "Latex" / "dissertation_figures" / "figure_reminder_workflow.png"

UI_DASHBOARD = ASSETS_DIR / "ui_dashboard.png"
UI_SCHEDULE = ASSETS_DIR / "ui_schedule.png"
UI_ASSISTANT = ASSETS_DIR / "ui_assistant.png"

FONT_NAME = "Century Gothic"

MIDNIGHT = RGBColor(18, 34, 62)
INDIGO = RGBColor(48, 63, 159)
TEAL = RGBColor(21, 128, 120)
AMBER = RGBColor(217, 119, 6)
ROSE = RGBColor(190, 24, 93)
SLATE = RGBColor(71, 85, 105)
MUTED = RGBColor(100, 116, 139)
INK = RGBColor(15, 23, 42)
WHITE = RGBColor(255, 255, 255)
PAPER = RGBColor(248, 250, 252)
SOFT_BLUE = RGBColor(224, 231, 255)
SOFT_TEAL = RGBColor(204, 251, 241)


def load_results() -> dict:
    with RESULTS_PATH.open("r", encoding="utf-8") as handle:
        return json.load(handle)


@lru_cache(maxsize=None)
def load_font(size: int, bold: bool = False):
    windows_font_dir = Path(os.environ.get("WINDIR", r"C:\Windows")) / "Fonts"
    candidates = [
        windows_font_dir / ("GOTHICB.TTF" if bold else "GOTHIC.TTF"),
        windows_font_dir / ("arialbd.ttf" if bold else "arial.ttf"),
    ]
    for candidate in candidates:
        if candidate.exists():
            return ImageFont.truetype(str(candidate), size=size)
    return ImageFont.load_default()


def vertical_gradient(size: tuple[int, int], top: tuple[int, int, int], bottom: tuple[int, int, int]) -> Image.Image:
    width, height = size
    image = Image.new("RGB", size, top)
    draw = ImageDraw.Draw(image)
    for y in range(height):
        mix = y / max(height - 1, 1)
        color = tuple(int(top[idx] + (bottom[idx] - top[idx]) * mix) for idx in range(3))
        draw.line((0, y, width, y), fill=color)
    return image


def rounded_panel(draw: ImageDraw.ImageDraw, box: tuple[int, int, int, int], fill: tuple[int, int, int],
                  outline: tuple[int, int, int] | None = None, radius: int = 28,
                  shadow: bool = True) -> None:
    x0, y0, x1, y1 = box
    if shadow:
        draw.rounded_rectangle((x0 + 12, y0 + 14, x1 + 12, y1 + 14), radius=radius, fill=(191, 219, 254))
    draw.rounded_rectangle(box, radius=radius, fill=fill, outline=outline, width=2 if outline else 0)


def text(draw: ImageDraw.ImageDraw, xy: tuple[int, int], value: str, size: int,
         fill: tuple[int, int, int], bold: bool = False) -> None:
    draw.text(xy, value, font=load_font(size, bold), fill=fill)


def pill(draw: ImageDraw.ImageDraw, xy: tuple[int, int], value: str,
         fill: tuple[int, int, int], fg: tuple[int, int, int], width: int | None = None) -> None:
    x, y = xy
    font = load_font(18, True)
    bbox = draw.textbbox((0, 0), value, font=font)
    pill_width = width or (bbox[2] - bbox[0] + 42)
    draw.rounded_rectangle((x, y, x + pill_width, y + 40), radius=20, fill=fill)
    draw.text((x + 21, y + 9), value, font=font, fill=fg)


def draw_browser_chrome(draw: ImageDraw.ImageDraw, box: tuple[int, int, int, int], title: str) -> None:
    x0, y0, x1, _ = box
    draw.rounded_rectangle((x0, y0, x1, y0 + 54), radius=24, fill=(243, 244, 246))
    for idx, color in enumerate(((239, 68, 68), (245, 158, 11), (34, 197, 94))):
        draw.ellipse((x0 + 24 + idx * 24, y0 + 18, x0 + 38 + idx * 24, y0 + 32), fill=color)
    text(draw, (x0 + 92, y0 + 15), title, 20, (71, 85, 105), bold=True)


def draw_header(draw: ImageDraw.ImageDraw, left: int, top: int, width: int, user_name: str) -> None:
    rounded_panel(draw, (left, top, left + width, top + 98), (255, 255, 255), shadow=False, radius=26)
    draw.rounded_rectangle((left + 16, top + 16, left + 80, top + 80), radius=22, fill=(79, 70, 229))
    text(draw, (left + 32, top + 34), "AI", 24, (255, 255, 255), bold=True)
    text(draw, (left + 104, top + 24), "Academic Task Manager", 28, (17, 24, 39), bold=True)
    text(draw, (left + 104, top + 57), "Intelligent task management powered by Agentic AI", 16, (100, 116, 139))

    tabs = [("Tasks", True), ("AI Schedule", False), ("AI Chat", False), ("Reminders", False)]
    tab_x = left + 630
    for label, active in tabs:
        fill = (79, 70, 229) if active else (238, 242, 255)
        fg = (255, 255, 255) if active else (79, 70, 229)
        pill(draw, (tab_x, top + 29), label, fill, fg)
        tab_x += 146 if label != "AI Schedule" else 162

    pill(draw, (left + width - 242, top + 29), user_name, (224, 231, 255), (55, 48, 163), width=142)
    pill(draw, (left + width - 88, top + 29), "Online", (220, 252, 231), (21, 128, 61), width=74)


def draw_stat_card(draw: ImageDraw.ImageDraw, box: tuple[int, int, int, int], title_value: str,
                   subtitle: str, accent: tuple[int, int, int]) -> None:
    rounded_panel(draw, box, (255, 255, 255), shadow=False, radius=24)
    x0, y0, _, _ = box
    draw.rounded_rectangle((x0 + 18, y0 + 18, x0 + 30, y0 + 74), radius=6, fill=accent)
    text(draw, (x0 + 52, y0 + 20), title_value, 30, (15, 23, 42), bold=True)
    text(draw, (x0 + 52, y0 + 58), subtitle, 16, (100, 116, 139))


def draw_field(draw: ImageDraw.ImageDraw, label: str, value: str, left: int, top: int, width: int,
               height: int = 54, multiline: bool = False) -> None:
    text(draw, (left, top - 26), label, 17, (71, 85, 105), bold=True)
    draw.rounded_rectangle((left, top, left + width, top + height), radius=16, fill=(248, 250, 252), outline=(226, 232, 240), width=2)
    value_y = top + 14 if not multiline else top + 12
    text(draw, (left + 18, value_y), value, 16, (148, 163, 184))


def create_dashboard_image(path: Path) -> None:
    image = vertical_gradient((1600, 900), (30, 27, 75), (79, 70, 229))
    draw = ImageDraw.Draw(image)
    draw.ellipse((1160, -80, 1620, 340), fill=(99, 102, 241))
    draw.ellipse((-140, 620, 360, 1080), fill=(67, 56, 202))

    window = (80, 56, 1520, 840)
    rounded_panel(draw, window, (245, 247, 251), radius=34)
    draw_browser_chrome(draw, window, "Prototype Frontend - Task Dashboard")
    draw_header(draw, 118, 130, 1364, "Amaka James")

    draw_stat_card(draw, (136, 248, 364, 336), "12", "Tracked tasks", (79, 70, 229))
    draw_stat_card(draw, (384, 248, 612, 336), "3", "Due this week", (13, 148, 136))
    draw_stat_card(draw, (632, 248, 860, 336), "1", "Urgent conflict", (217, 119, 6))
    draw_stat_card(draw, (880, 248, 1108, 336), "9.2s", "Avg task entry", (190, 24, 93))

    rounded_panel(draw, (136, 364, 566, 786), (255, 255, 255), radius=28)
    text(draw, (164, 392), "Add New Task", 28, (15, 23, 42), bold=True)
    text(draw, (164, 428), "Structured input form with validation and UX timing metric", 16, (100, 116, 139))
    draw_field(draw, "Task Title", "Complete Chapter 4 evaluation write-up", 164, 488, 374)
    draw_field(draw, "Description", "Compare baseline and agent outcomes, then draft results summary.", 164, 578, 374, height=92, multiline=True)
    draw_field(draw, "Deadline", "22 Apr 2026", 164, 702, 176)
    draw_field(draw, "Priority", "High", 362, 702, 176)
    draw.rounded_rectangle((164, 756, 538, 804), radius=18, fill=(16, 185, 129))
    text(draw, (286, 769), "Add Task", 20, (255, 255, 255), bold=True)

    rounded_panel(draw, (598, 364, 1466, 786), (255, 255, 255), radius=28)
    text(draw, (626, 392), "Tasks for Amaka James", 28, (15, 23, 42), bold=True)
    pill(draw, (1246, 389), "12 tasks", (79, 70, 229), (255, 255, 255))
    draw_field(draw, "Search", "Search title or description", 626, 470, 258)
    draw_field(draw, "Status", "Pending", 904, 470, 124)
    draw_field(draw, "Priority", "All", 1048, 470, 124)
    draw_field(draw, "Sort", "Deadline", 1192, 470, 124)

    task_cards = [
        ("Run evaluation experiments", "Due 18 Apr 2026", "High", "Pending", (254, 226, 226), (185, 28, 28)),
        ("Refine frontend screenshots", "Due 19 Apr 2026", "Medium", "In Progress", (224, 231, 255), (55, 48, 163)),
        ("Finalize dissertation slides", "Due 21 Apr 2026", "High", "Pending", (254, 226, 226), (185, 28, 28)),
        ("Check reminder cron logs", "Due 24 Apr 2026", "Low", "Completed", (220, 252, 231), (21, 128, 61)),
    ]
    y = 550
    for title_value, due_text, priority, status, status_fill, status_fg in task_cards:
        rounded_panel(draw, (626, y, 1438, y + 72), (248, 250, 252), shadow=False, radius=22)
        text(draw, (648, y + 14), title_value, 21, (15, 23, 42), bold=True)
        text(draw, (648, y + 42), due_text, 15, (100, 116, 139))
        pill(draw, (1132, y + 16), priority, (255, 247, 237), (194, 65, 12), width=94)
        pill(draw, (1244, y + 16), status, status_fill, status_fg, width=150)
        y += 86

    draw.rounded_rectangle((626, 744, 1438, 786), radius=18, fill=(255, 247, 237))
    text(draw, (648, 756), "Conflict detected: Dissertation slides and evaluation experiments compete for the same locked study block.", 15, (154, 52, 18))
    path.parent.mkdir(parents=True, exist_ok=True)
    image.save(path)


def create_schedule_image(path: Path) -> None:
    image = vertical_gradient((1600, 900), (15, 23, 42), (30, 41, 59))
    draw = ImageDraw.Draw(image)
    draw.ellipse((1180, -120, 1640, 280), fill=(14, 165, 233))
    draw.ellipse((-120, 690, 360, 1080), fill=(45, 212, 191))

    window = (80, 56, 1520, 840)
    rounded_panel(draw, window, (246, 248, 252), radius=34)
    draw_browser_chrome(draw, window, "Prototype Frontend - AI Schedule")
    draw_header(draw, 118, 130, 1364, "Amaka James")

    rounded_panel(draw, (136, 248, 502, 786), (255, 255, 255), radius=28)
    text(draw, (164, 278), "Your Availability", 28, (15, 23, 42), bold=True)
    text(draw, (164, 314), "Study slots are stored and reused during schedule generation", 16, (100, 116, 139))
    draw_field(draw, "Day", "Tuesday", 164, 390, 148)
    draw_field(draw, "From", "18:00", 332, 390, 92)
    draw_field(draw, "To", "21:00", 436, 390, 92)
    draw.rounded_rectangle((164, 462, 474, 508), radius=18, fill=(79, 70, 229))
    text(draw, (290, 474), "Add Slot", 20, (255, 255, 255), bold=True)

    slots = [
        ("Monday", "18:00 - 21:00", "Library"),
        ("Tuesday", "18:00 - 21:00", "Home desk"),
        ("Thursday", "17:30 - 20:30", "Lab"),
        ("Saturday", "09:00 - 13:00", "Quiet room"),
    ]
    y = 548
    for day, time_range, location in slots:
        rounded_panel(draw, (164, y, 474, y + 62), (248, 250, 252), shadow=False, radius=20)
        text(draw, (184, y + 12), day, 18, (30, 41, 59), bold=True)
        text(draw, (184, y + 34), time_range, 15, (100, 116, 139))
        text(draw, (336, y + 22), location, 15, (79, 70, 229))
        y += 74

    rounded_panel(draw, (538, 248, 1466, 786), (255, 255, 255), radius=28)
    text(draw, (566, 278), "AI-Generated Schedule", 28, (15, 23, 42), bold=True)
    text(draw, (566, 314), "The scheduler places tasks into free slots and explains difficult trade-offs", 16, (100, 116, 139))
    draw.rounded_rectangle((1224, 274, 1438, 320), radius=18, fill=(79, 70, 229))
    text(draw, (1270, 286), "Generate Schedule", 18, (255, 255, 255), bold=True)

    draw.rounded_rectangle((566, 364, 1438, 420), radius=20, fill=(236, 254, 255))
    text(draw, (590, 378), "AI reasoning: deadlines and locked blocks were balanced first, then workload was spread across remaining study windows.", 16, (15, 118, 110))

    day_boxes = [
        ("Tue 16 Apr", 572, (79, 70, 229)),
        ("Thu 18 Apr", 864, (13, 148, 136)),
        ("Sat 20 Apr", 1156, (217, 119, 6)),
    ]
    for label, x, accent in day_boxes:
        rounded_panel(draw, (x, 454, x + 248, 764), (248, 250, 252), shadow=False, radius=24)
        draw.rounded_rectangle((x + 18, 474, x + 230, 514), radius=16, fill=accent)
        text(draw, (x + 52, 484), label, 18, (255, 255, 255), bold=True)

    schedule_content = [
        [
            ("18:00 - 19:20", "Prepare evaluation tables", (224, 231, 255)),
            ("19:30 - 20:45", "Revise Chapter 5 narrative", (219, 234, 254)),
        ],
        [
            ("17:30 - 18:45", "Resolve reminder SQL issue", (204, 251, 241)),
            ("19:00 - 20:20", "Produce results chart", (209, 250, 229)),
        ],
        [
            ("09:00 - 10:40", "Refine presentation deck", (254, 243, 199)),
            ("11:00 - 12:30", "Proofread conclusion slide", (254, 240, 138)),
        ],
    ]
    for col_idx, entries in enumerate(schedule_content):
        x = day_boxes[col_idx][1] + 18
        y = 536
        for time_range, task_title, fill in entries:
            rounded_panel(draw, (x, y, x + 212, y + 92), fill, shadow=False, radius=20)
            text(draw, (x + 18, y + 14), time_range, 16, (71, 85, 105), bold=True)
            text(draw, (x + 18, y + 44), task_title, 19, (15, 23, 42), bold=True)
            y += 108

    draw.rounded_rectangle((566, 730, 1438, 776), radius=18, fill=(220, 252, 231))
    text(draw, (590, 744), "Schedule generated successfully. No hard deadline misses in this 7-day planning window.", 16, (21, 128, 61))
    path.parent.mkdir(parents=True, exist_ok=True)
    image.save(path)


def create_assistant_image(path: Path) -> None:
    image = vertical_gradient((1600, 900), (17, 24, 39), (67, 56, 202))
    draw = ImageDraw.Draw(image)
    draw.ellipse((1200, -120, 1650, 320), fill=(129, 140, 248))
    draw.ellipse((-180, 650, 360, 1100), fill=(45, 212, 191))

    window = (80, 56, 1520, 840)
    rounded_panel(draw, window, (245, 247, 251), radius=34)
    draw_browser_chrome(draw, window, "Prototype Frontend - Reminders and AI Chat")
    draw_header(draw, 118, 130, 1364, "Amaka James")

    rounded_panel(draw, (136, 248, 580, 786), (255, 255, 255), radius=28)
    text(draw, (164, 278), "Send Task Reminder", 28, (15, 23, 42), bold=True)
    text(draw, (164, 314), "Manual reminder trigger mirrors the hosted Supabase Cron pathway", 16, (100, 116, 139))
    draw.rounded_rectangle((164, 362, 548, 438), radius=20, fill=(254, 249, 195))
    text(draw, (190, 382), "Upcoming deadlines: 3 due soon, 1 overdue", 18, (133, 77, 14), bold=True)
    draw_field(draw, "Reminder Period", "Tasks due in 5 days", 164, 492, 344)
    draw_field(draw, "Recipient", "amaka.james@example.com", 164, 584, 344)
    draw.rounded_rectangle((164, 676, 508, 726), radius=18, fill=(217, 119, 6))
    text(draw, (260, 690), "Send Reminder Email", 20, (255, 255, 255), bold=True)

    rounded_panel(draw, (164, 742, 548, 796), (220, 252, 231), shadow=False, radius=18)
    text(draw, (188, 758), "Delivery status: reminder queued through Brevo without UI errors.", 15, (21, 128, 61))

    rounded_panel(draw, (618, 248, 1466, 786), (255, 255, 255), radius=28)
    text(draw, (646, 278), "AI Task Assistant", 28, (15, 23, 42), bold=True)
    text(draw, (646, 314), "Natural-language commands for task review, task creation, and lightweight assistant guidance", 16, (100, 116, 139))

    conversations = [
        ("user", "What tasks are due this week?"),
        ("assistant", "You have 3 tasks due within 7 days. The most urgent are evaluation experiments and the slide deck."),
        ("user", "Add a task to rehearse the dissertation presentation on Saturday morning."),
        ("assistant", "Task added. I also recommend keeping 30 minutes on Sunday evening for final slide edits."),
    ]
    y = 376
    for role, message in conversations:
        bubble_fill = (224, 231, 255) if role == "user" else (241, 245, 249)
        bubble_x0 = 912 if role == "user" else 646
        bubble_x1 = 1428 if role == "user" else 1218
        rounded_panel(draw, (bubble_x0, y, bubble_x1, y + 84), bubble_fill, shadow=False, radius=22)
        label = "You" if role == "user" else "Assistant"
        text(draw, (bubble_x0 + 22, y + 14), label, 16, (79, 70, 229) if role == "user" else (15, 23, 42), bold=True)
        text(draw, (bubble_x0 + 22, y + 40), message, 16, (30, 41, 59))
        y += 102

    draw.rounded_rectangle((646, 718, 1428, 772), radius=18, fill=(248, 250, 252), outline=(226, 232, 240), width=2)
    text(draw, (672, 734), "Type your message... e.g. 'show high priority tasks'", 16, (148, 163, 184))
    draw.rounded_rectangle((1304, 722, 1408, 768), radius=16, fill=(79, 70, 229))
    text(draw, (1340, 734), "Send", 18, (255, 255, 255), bold=True)

    draw.rounded_rectangle((618, 804, 1466, 836), radius=16, fill=(224, 231, 255))
    text(draw, (646, 812), "Offline support: queued actions sync automatically when the client reconnects.", 14, (55, 48, 163))
    path.parent.mkdir(parents=True, exist_ok=True)
    image.save(path)


def ensure_assets() -> None:
    ASSETS_DIR.mkdir(parents=True, exist_ok=True)
    create_dashboard_image(UI_DASHBOARD)
    create_schedule_image(UI_SCHEDULE)
    create_assistant_image(UI_ASSISTANT)


def set_background(slide, color: RGBColor = PAPER) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_soft_shapes(slide) -> None:
    circle_top = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(10.8), Inches(-0.8), Inches(3.2), Inches(3.2))
    circle_top.fill.solid()
    circle_top.fill.fore_color.rgb = SOFT_BLUE
    circle_top.line.fill.background()
    circle_bottom = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(-0.7), Inches(6.0), Inches(2.6), Inches(2.6))
    circle_bottom.fill.solid()
    circle_bottom.fill.fore_color.rgb = SOFT_TEAL
    circle_bottom.line.fill.background()


def add_slide_header(slide, section: str, title_value: str, subtitle: str | None = None) -> None:
    add_soft_shapes(slide)
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.18))
    band.fill.solid()
    band.fill.fore_color.rgb = MIDNIGHT
    band.line.fill.background()

    section_box = slide.shapes.add_textbox(Inches(0.72), Inches(0.36), Inches(3.0), Inches(0.26))
    p = section_box.text_frame.paragraphs[0]
    p.text = section.upper()
    p.font.name = FONT_NAME
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = TEAL

    title_box = slide.shapes.add_textbox(Inches(0.72), Inches(0.68), Inches(9.6), Inches(0.6))
    p = title_box.text_frame.paragraphs[0]
    p.text = title_value
    p.font.name = FONT_NAME
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = INK

    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.72), Inches(1.12), Inches(10.8), Inches(0.36))
        p = subtitle_box.text_frame.paragraphs[0]
        p.text = subtitle
        p.font.name = FONT_NAME
        p.font.size = Pt(13)
        p.font.color.rgb = SLATE


def add_footer(slide, page_number: int) -> None:
    footer = slide.shapes.add_textbox(Inches(0.72), Inches(7.05), Inches(11.9), Inches(0.22))
    p = footer.text_frame.paragraphs[0]
    p.text = "Agentic AI for Education  |  MSc Dissertation Presentation"
    p.font.name = FONT_NAME
    p.font.size = Pt(10)
    p.font.color.rgb = MUTED
    p.alignment = PP_ALIGN.LEFT

    page = slide.shapes.add_textbox(Inches(12.1), Inches(7.0), Inches(0.5), Inches(0.25))
    p = page.text_frame.paragraphs[0]
    p.text = str(page_number)
    p.font.name = FONT_NAME
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = MIDNIGHT
    p.alignment = PP_ALIGN.RIGHT


def add_text(slide, left: float, top: float, width: float, height: float, value: str,
             size: int = 18, bold: bool = False, color: RGBColor = INK,
             align: PP_ALIGN = PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    p = box.text_frame.paragraphs[0]
    p.text = value
    p.font.name = FONT_NAME
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    box.text_frame.word_wrap = True
    return box


def add_bullets(slide, items: Iterable[str], left: float, top: float, width: float, height: float,
                font_size: int = 18, color: RGBColor = INK):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = f"• {item}"
        p.font.name = FONT_NAME
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.level = 0
        p.space_after = Pt(10)
    return box


def add_card(slide, left: float, top: float, width: float, height: float, title_value: str,
             body: str, accent: RGBColor = TEAL, fill: RGBColor = WHITE,
             body_size: int = 15, title_size: int = 17) -> None:
    shadow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left + 0.05), Inches(top + 0.06), Inches(width), Inches(height))
    shadow.fill.solid()
    shadow.fill.fore_color.rgb = SOFT_BLUE
    shadow.line.fill.background()

    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    card.line.color.rgb = accent
    card.line.width = Pt(1.4)

    accent_bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(left + 0.16), Inches(top + 0.18), Inches(0.1), Inches(height - 0.36))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = accent
    accent_bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(left + 0.34), Inches(top + 0.16), Inches(width - 0.46), Inches(0.34))
    p = title_box.text_frame.paragraphs[0]
    p.text = title_value
    p.font.name = FONT_NAME
    p.font.size = Pt(title_size)
    p.font.bold = True
    p.font.color.rgb = INK

    body_box = slide.shapes.add_textbox(Inches(left + 0.34), Inches(top + 0.54), Inches(width - 0.5), Inches(height - 0.68))
    p = body_box.text_frame.paragraphs[0]
    p.text = body
    p.font.name = FONT_NAME
    p.font.size = Pt(body_size)
    p.font.color.rgb = SLATE
    body_box.text_frame.word_wrap = True


def add_image_frame(slide, image_path: Path, left: float, top: float, width: float, height: float | None = None) -> None:
    shadow_height = height or width * 0.56
    shadow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left + 0.06), Inches(top + 0.08), Inches(width), Inches(shadow_height))
    shadow.fill.solid()
    shadow.fill.fore_color.rgb = SOFT_BLUE
    shadow.line.fill.background()
    kwargs = {"width": Inches(width)}
    if height is not None:
        kwargs["height"] = Inches(height)
    slide.shapes.add_picture(str(image_path), Inches(left), Inches(top), **kwargs)


def title_slide(prs: Presentation, page_number: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, MIDNIGHT)

    halo_one = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(9.6), Inches(-0.7), Inches(4.5), Inches(4.5))
    halo_one.fill.solid()
    halo_one.fill.fore_color.rgb = INDIGO
    halo_one.line.fill.background()
    halo_two = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(-0.6), Inches(5.2), Inches(3.4), Inches(3.4))
    halo_two.fill.solid()
    halo_two.fill.fore_color.rgb = TEAL
    halo_two.line.fill.background()

    tag = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.72), Inches(0.62), Inches(2.7), Inches(0.38))
    tag.fill.solid()
    tag.fill.fore_color.rgb = SOFT_TEAL
    tag.line.fill.background()
    add_text(slide, 0.95, 0.72, 2.2, 0.18, "MSC DISSERTATION PROJECT", size=11, bold=True, color=TEAL)

    add_text(slide, 0.74, 1.42, 5.7, 0.75, "Agentic AI for Education", size=30, bold=True, color=WHITE)
    add_text(slide, 0.74, 2.18, 5.7, 0.72, "Design, implementation and evaluation of a learner-support system", size=19, color=SOFT_BLUE)
    add_bullets(
        slide,
        [
            "A working prototype that supports task capture, AI scheduling, replanning, reminders, and natural-language interaction.",
            "Evaluation framed around stable scheduling, dynamic change, and severe workload saturation rather than inflated claims.",
            "Presentation refreshed with interface visuals and a more dissertation-appropriate narrative structure.",
        ],
        left=0.92,
        top=3.0,
        width=5.6,
        height=2.3,
        font_size=18,
        color=WHITE,
    )

    add_card(slide, 0.92, 5.7, 2.3, 0.8, "Stack", "React PWA, Flask, PostgreSQL, LangGraph, Supabase, Brevo", accent=AMBER, fill=WHITE, body_size=14)
    add_card(slide, 3.38, 5.7, 2.3, 0.8, "Claim", "Active learner support is demonstrated, but optimisation limits remain visible.", accent=TEAL, fill=WHITE, body_size=14)

    add_image_frame(slide, UI_DASHBOARD, 7.08, 1.1, 5.3)
    add_image_frame(slide, UI_SCHEDULE, 8.1, 4.15, 4.15)
    add_image_frame(slide, UI_ASSISTANT, 10.0, 3.3, 2.25)
    add_footer(slide, page_number)


def problem_slide(prs: Presentation, page_number: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_slide_header(slide, "Research Framing", "Problem, aim and research questions", "Why passive study tools are not enough for a dissertation-level systems project")
    add_bullets(
        slide,
        [
            "Most academic productivity tools record tasks, but they do not actively schedule work, replan after disruption, or sustain follow-through.",
            "This project asks whether an agentic learner-support system can deliver dependable scheduling, adaptation, explanation, and proactive reminders.",
            "The dissertation contribution is the system as a whole: frontend, backend, orchestration logic, evaluation framework, and hosted reminder workflow.",
        ],
        left=0.84,
        top=1.72,
        width=6.1,
        height=3.25,
    )
    add_card(slide, 7.28, 1.68, 5.2, 1.26, "Aim", "Design and evaluate a learner-support system that does more than store deadlines: it should actively plan, adapt, and remind.", accent=TEAL, fill=WHITE, body_size=16, title_size=18)
    add_card(slide, 7.28, 3.12, 5.2, 1.52, "Research questions", "RQ1: Can the system generate viable schedules?  RQ2: Does agentic behaviour improve schedule quality over a baseline?  RQ3: Can it adapt credibly when tasks or time windows change?", accent=AMBER, fill=WHITE, body_size=16, title_size=18)
    add_card(slide, 7.28, 4.88, 5.2, 1.1, "Why it is MSc-worthy", "The work combines implementation, system integration, debugging, evaluation, and a bounded research claim rather than a simple CRUD application.", accent=ROSE, fill=WHITE, body_size=16, title_size=18)
    add_footer(slide, page_number)


def contributions_slide(prs: Presentation, page_number: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_slide_header(slide, "Prototype Scope", "What was actually built", "Core features are framed as system contributions, not marketing claims")
    add_card(slide, 0.82, 1.72, 2.8, 1.2, "Task capture", "Structured forms, validation, deadline tracking, filterable task views, and offline queue awareness.", accent=INDIGO, fill=WHITE, body_size=15)
    add_card(slide, 3.84, 1.72, 2.8, 1.2, "AI scheduling", "Availability-aware schedule generation with reasoning, conflict visibility, and dynamic replanning triggers.", accent=TEAL, fill=WHITE, body_size=15)
    add_card(slide, 0.82, 3.12, 2.8, 1.2, "Natural-language support", "Chat commands for task review, task creation, and lightweight assistant guidance inside the prototype.", accent=AMBER, fill=WHITE, body_size=15)
    add_card(slide, 3.84, 3.12, 2.8, 1.2, "Hosted reminders", "Supabase Cron invokes an Edge Function that queries due tasks and delivers email through Brevo.", accent=ROSE, fill=WHITE, body_size=15)
    add_text(slide, 0.84, 4.72, 5.9, 0.42, "These features matter because they expose agentic behaviour in concrete user-facing workflows.", size=16, color=SLATE)
    add_image_frame(slide, USE_CASE_PATH, 6.95, 1.72, 5.0)
    add_text(slide, 7.02, 5.08, 4.9, 0.42, "Use-case coverage keeps the talk grounded in what the system supports, not only how it is implemented.", size=15, color=SLATE)
    add_footer(slide, page_number)


def dashboard_ui_slide(prs: Presentation, page_number: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_slide_header(slide, "Frontend Evidence", "Prototype UI: task management dashboard", "Adding interface visuals makes the deck substantially more credible for a dissertation presentation")
    add_image_frame(slide, UI_DASHBOARD, 0.82, 1.7, 8.2)
    add_card(slide, 9.28, 1.78, 3.1, 1.06, "Why keep this slide", "Examiners can see that the work is implemented, not only described abstractly.", accent=TEAL, body_size=15)
    add_card(slide, 9.28, 3.0, 3.1, 1.18, "What the image shows", "Task creation, filtering, priority cues, conflict surfacing, and user-state indicators.", accent=AMBER, body_size=15)
    add_card(slide, 9.28, 4.34, 3.1, 1.3, "Presentation tip", "Talk through one realistic workflow: create task, inspect priority, then show how the system surfaces a scheduling issue.", accent=ROSE, body_size=15)
    add_footer(slide, page_number)


def advanced_ui_slide(prs: Presentation, page_number: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_slide_header(slide, "Frontend Evidence", "Prototype UI: scheduling, reminders and assistant support", "These screens connect the UI to the dissertation claim about active learner support")
    add_image_frame(slide, UI_SCHEDULE, 0.82, 1.74, 5.98)
    add_image_frame(slide, UI_ASSISTANT, 6.54, 1.74, 5.98)
    add_text(slide, 0.92, 5.44, 5.7, 0.32, "Left: availability-aware schedule generation with reasoning and visible day-by-day slots.", size=15, color=SLATE)
    add_text(slide, 6.64, 5.44, 5.7, 0.32, "Right: reminder dispatch controls, assistant conversation, and offline-safe UX cues.", size=15, color=SLATE)
    add_card(slide, 0.92, 5.92, 3.45, 0.82, "Why it helps in the viva", "It demonstrates that the system supports planning, follow-through, and intervention, not just record-keeping.", accent=TEAL, body_size=14)
    add_card(slide, 4.58, 5.92, 3.45, 0.82, "Bounded claim", "These are prototype screenshots. The argument is implementation credibility, not production polish.", accent=AMBER, body_size=14)
    add_card(slide, 8.24, 5.92, 3.45, 0.82, "Key message", "The frontend operationalises the research idea so the agentic behaviour is visible and testable.", accent=INDIGO, body_size=14)
    add_footer(slide, page_number)


def architecture_slide(prs: Presentation, page_number: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_slide_header(slide, "System Design", "Architecture and integration story", "Explain the software boundary clearly: client, backend logic, orchestration, and hosted services")
    add_image_frame(slide, ARCH_PATH, 0.82, 1.78, 7.2)
    add_card(slide, 8.38, 1.9, 3.8, 1.12, "Frontend", "React PWA manages user interaction, session state, offline queue feedback, and feature tabs.", accent=INDIGO, body_size=15)
    add_card(slide, 8.38, 3.18, 3.8, 1.12, "Backend", "Flask routes requests to task operations, schedule generation, evaluation support, and conflict analysis.", accent=TEAL, body_size=15)
    add_card(slide, 8.38, 4.46, 3.8, 1.12, "Hosted path", "Supabase stores operational data and executes the deadline-reminder Edge Function on a daily cron.", accent=AMBER, body_size=15)
    add_text(slide, 0.9, 5.98, 7.0, 0.44, "The architecture slide is stronger when you talk about why each boundary exists, not when you simply list technologies.", size=15, color=SLATE)
    add_footer(slide, page_number)


def workflow_slide(prs: Presentation, page_number: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_slide_header(slide, "System Operation", "Operational workflow and hosted reminder automation", "Two diagrams are enough here: one for request flow and one for proactive intervention")
    add_image_frame(slide, REQUEST_PATH, 0.82, 1.8, 5.78)
    add_image_frame(slide, REMINDER_PATH, 6.74, 1.8, 5.78)
    add_text(slide, 0.92, 5.34, 5.46, 0.3, "Request workflow: task input, scheduling logic, conflict analysis, and response generation.", size=15, color=SLATE)
    add_text(slide, 6.84, 5.34, 5.46, 0.3, "Reminder workflow: Supabase Cron, Edge Function execution, overdue detection, and Brevo delivery.", size=15, color=SLATE)
    add_card(slide, 0.92, 5.84, 11.4, 0.8, "Presentation tip", "Use this slide to show that the reminder capability is not hypothetical. It exists as a hosted path with concrete failure diagnosis and logging.", accent=ROSE, body_size=14, title_size=16)
    add_footer(slide, page_number)


def evaluation_design_slide(prs: Presentation, page_number: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_slide_header(slide, "Evaluation", "Evaluation design and scenario coverage", "The deck should show a disciplined evaluation method, not only screenshots and architecture")
    rows = [
        ("A", "Standard week", "Baseline scheduling quality under ordinary workload"),
        ("B", "Deadline compression", "Behaviour under clustered urgent deadlines"),
        ("C", "Disruption", "Adaptation after time blocks are removed"),
        ("D", "Dynamic addition", "Replanning when urgent new tasks arrive"),
        ("E", "Five-week project", "Extended workload with visible saturation risk"),
    ]
    table = slide.shapes.add_table(len(rows) + 1, 3, Inches(0.82), Inches(1.84), Inches(7.7), Inches(3.6)).table
    table.columns[0].width = Inches(0.8)
    table.columns[1].width = Inches(2.0)
    table.columns[2].width = Inches(4.9)
    headers = ["Scenario", "Simulation", "Purpose"]
    for idx, header in enumerate(headers):
        cell = table.cell(0, idx)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = MIDNIGHT
        for p in cell.text_frame.paragraphs:
            p.font.name = FONT_NAME
            p.font.bold = True
            p.font.size = Pt(13)
            p.font.color.rgb = WHITE
    for row_idx, row in enumerate(rows, start=1):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.text = value
            cell.fill.solid()
            cell.fill.fore_color.rgb = PAPER if row_idx % 2 else WHITE
            for p in cell.text_frame.paragraphs:
                p.font.name = FONT_NAME
                p.font.size = Pt(12)
                p.font.color.rgb = INK

    add_card(slide, 8.9, 1.9, 3.3, 1.08, "Metrics", "Generation time, conflict rate, deadline compliance, workload spread, adaptation speed, and efficiency.", accent=TEAL, body_size=15)
    add_card(slide, 8.9, 3.16, 3.3, 1.08, "Comparator", "A rule-based earliest-deadline-first baseline scheduler provides a defensible reference point.", accent=AMBER, body_size=15)
    add_card(slide, 8.9, 4.42, 3.3, 1.08, "Why it matters", "This evaluation design supports a bounded academic claim rather than vague statements about intelligence.", accent=ROSE, body_size=15)
    add_footer(slide, page_number)


def results_slide(prs: Presentation, results: dict, page_number: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_slide_header(slide, "Evaluation", "Results and dissertation claim", "The key is to present evidence honestly: gains in some areas, limits in others")

    chart_data = CategoryChartData()
    chart_data.categories = ["A", "B", "D", "E"]
    chart_data.add_series(
        "Baseline conflict %",
        [
            results["scenario_a"]["baseline"]["avg_conflict_rate"],
            results["scenario_b"]["baseline"]["avg_conflict_rate"],
            results["scenario_d"]["baseline"]["avg_conflict_rate"],
            results["scenario_e"]["baseline"]["avg_conflict_rate"],
        ],
    )
    chart_data.add_series(
        "Agent conflict %",
        [
            results["scenario_a"]["agent"]["avg_conflict_rate"],
            results["scenario_b"]["agent"]["avg_conflict_rate"],
            results["scenario_d"]["agent"]["avg_conflict_rate"],
            results["scenario_e"]["agent"]["avg_conflict_rate"],
        ],
    )
    chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.82), Inches(1.9), Inches(6.8), Inches(3.5), chart_data).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.value_axis.maximum_scale = 100
    chart.value_axis.minimum_scale = 0
    chart.value_axis.tick_labels.font.size = Pt(11)
    chart.category_axis.tick_labels.font.size = Pt(12)
    for idx, series in enumerate(chart.series):
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = INDIGO if idx == 0 else TEAL

    summary = results["summary"]
    h1 = "Supported" if summary["hypothesis_h1_supported"] else "Not supported"
    h2 = "Supported" if summary["hypothesis_h2_supported"] else "Not supported"
    h3 = "Supported" if summary["hypothesis_h3_supported"] else "Not supported"
    add_card(slide, 8.0, 1.92, 4.28, 0.98, f"H1 speed: {h1}", "The agent is slower than the baseline, so the dissertation should not claim raw performance superiority.", accent=ROSE, body_size=14)
    add_card(slide, 8.0, 3.04, 4.28, 0.98, f"H2 schedule quality: {h2}", "Conflict handling improves modestly in the better scenarios, which supports a limited quality claim.", accent=TEAL, body_size=14)
    add_card(slide, 8.0, 4.16, 4.28, 0.98, f"H3 adaptation: {h3}", "Dynamic-change scenarios provide the most defensible evidence for agentic behaviour.", accent=AMBER, body_size=14)

    a_agent = results["scenario_a"]["agent"]["avg_conflict_rate"]
    a_base = results["scenario_a"]["baseline"]["avg_conflict_rate"]
    d_agent = results["scenario_d"]["agent"]["avg_conflict_rate"]
    d_base = results["scenario_d"]["baseline"]["avg_conflict_rate"]
    e_comp = results["scenario_e"]["agent"]["avg_deadline_compliance_rate"]
    add_text(slide, 0.92, 5.72, 3.7, 0.34, f"Scenario A conflict rate: {a_base:.1f}% baseline vs {a_agent:.1f}% agent", size=15, color=SLATE)
    add_text(slide, 0.92, 6.04, 3.7, 0.34, f"Scenario D conflict rate: {d_base:.1f}% baseline vs {d_agent:.1f}% agent", size=15, color=SLATE)
    add_card(slide, 4.72, 5.68, 7.56, 0.86, "Defensible conclusion", f"The prototype demonstrates active learner support, especially for replanning and user-facing intervention. Severe overload remains unresolved, with Scenario E compliance at {e_comp:.1f}%.", accent=INDIGO, body_size=14)
    add_footer(slide, page_number)


def reflection_slide(prs: Presentation, page_number: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_slide_header(slide, "Reflection", "Limitations, critical reflection and next steps", "This slide matters because dissertation presentations are judged partly on how honestly they handle limits")
    add_card(slide, 0.86, 1.84, 3.65, 1.28, "Current limitations", "The prototype is not a production planner. Scenario E exposes feasibility limits, and reminder logic depends on keeping hosted SQL aligned with intended behaviour.", accent=ROSE, body_size=15)
    add_card(slide, 4.84, 1.84, 3.65, 1.28, "Method limits", "The evaluation is simulation-based. A stronger future study would include real students, perceived usefulness, and sustained usage data.", accent=AMBER, body_size=15)
    add_card(slide, 8.82, 1.84, 3.65, 1.28, "Engineering reflection", "Much of the value lies in integration, debugging, deployment, and test coverage across several moving parts.", accent=TEAL, body_size=15)
    add_bullets(
        slide,
        [
            "Add task decomposition and workload-feasibility warnings before trying to optimise every scenario.",
            "Expand observability around reminder dispatch, failure logging, and hosted environment drift.",
            "Refine interface polish later; the immediate academic value is in demonstrable behaviour and evaluation evidence.",
        ],
        left=0.92,
        top=3.64,
        width=11.3,
        height=1.8,
    )
    add_card(slide, 0.92, 5.66, 11.4, 0.9, "Why the deck now works better", "It no longer looks like a generic software demo. It now foregrounds research framing, interface evidence, architecture, evaluation, limitations, and a bounded conclusion.", accent=INDIGO, body_size=15, title_size=16)
    add_footer(slide, page_number)


def closing_slide(prs: Presentation, page_number: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, MIDNIGHT)
    halo = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(9.8), Inches(-0.6), Inches(4.2), Inches(4.2))
    halo.fill.solid()
    halo.fill.fore_color.rgb = INDIGO
    halo.line.fill.background()
    add_text(slide, 0.82, 1.1, 5.6, 0.4, "Closing Position", size=12, bold=True, color=SOFT_TEAL)
    add_text(slide, 0.82, 1.55, 7.2, 0.8, "A credible MSc dissertation presentation should show both system substance and critical judgement.", size=28, bold=True, color=WHITE)
    add_bullets(
        slide,
        [
            "The project produced a working learner-support prototype with scheduling, adaptation, reminders, and assistant features.",
            "The evidence is strongest in implementation depth, system integration, and bounded evaluation results.",
            "The right final claim is not 'perfect planner' but 'implemented and evaluated agentic learner-support system with visible strengths and visible limits'.",
        ],
        left=0.96,
        top=2.72,
        width=6.0,
        height=2.4,
        font_size=18,
        color=WHITE,
    )
    add_card(slide, 7.3, 2.1, 4.6, 1.15, "Answer to your question", "Before this refresh: not quite. With UI visuals, cleaner structure, Century Gothic body text, and stronger academic framing: much closer to something that can fly.", accent=AMBER, fill=WHITE, body_size=15)
    add_card(slide, 7.3, 3.55, 4.6, 1.05, "What to rehearse", "Explain the contribution, then show one UI workflow, one architecture slide, one evaluation slide, and one limitation slide.", accent=TEAL, fill=WHITE, body_size=15)
    add_text(slide, 0.96, 6.3, 3.8, 0.4, "Questions", size=24, bold=True, color=SOFT_BLUE)
    add_footer(slide, page_number)


def write_notes() -> None:
    notes = """# Agentic AI for Education Presentation Notes

1. Title: present the work as a learner-support system, not just an AI feature demo.
2. Problem, aim and research questions: explain the gap between passive task tracking and active learner support.
3. Prototype scope: walk through the four implemented contribution areas and keep the claims concrete.
4. Task dashboard UI: use the screenshot to prove implementation depth and describe one realistic user workflow.
5. Scheduling, reminders and assistant UI: connect each screen directly to the dissertation claim about agentic behaviour.
6. Architecture and integration story: explain why the frontend, backend, orchestration logic, and hosted services are separated.
7. Operational workflow and reminder automation: show that proactive intervention is implemented as a real hosted path.
8. Evaluation design: highlight scenarios, metrics, and the baseline comparison.
9. Results and dissertation claim: stress supported hypotheses, reject the speed claim, and keep the conclusion bounded.
10. Reflection: show awareness of overload limits, simulation-only evaluation, and production-readiness gaps.
11. Closing: position the project as an implemented and evaluated prototype with credible strengths and honest limits.
"""
    NOTES_PATH.write_text(notes, encoding="utf-8")


def build_presentation() -> None:
    ensure_assets()
    results = load_results()
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    title_slide(prs, 1)
    problem_slide(prs, 2)
    contributions_slide(prs, 3)
    dashboard_ui_slide(prs, 4)
    advanced_ui_slide(prs, 5)
    architecture_slide(prs, 6)
    workflow_slide(prs, 7)
    evaluation_design_slide(prs, 8)
    results_slide(prs, results, 9)
    reflection_slide(prs, 10)
    closing_slide(prs, 11)

    prs.save(str(OUTPUT_PATH))
    write_notes()
    if PROTOTYPE_OUTPUT_PATH.parent.exists():
        PROTOTYPE_OUTPUT_PATH.write_bytes(OUTPUT_PATH.read_bytes())


if __name__ == "__main__":
    build_presentation()
