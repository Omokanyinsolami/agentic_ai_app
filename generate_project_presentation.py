from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from Backend.evaluation import TestRunner


ROOT = Path(__file__).resolve().parent
OUTPUT = ROOT / "Agentic_AI_for_Education_Presentation.pptx"
NOTES = ROOT / "Agentic_AI_for_Education_Presentation_Notes.md"
ARCH = ROOT / "Latex" / "prototype_architecture.png"
USE_CASE = ROOT / "Latex" / "dissertation_figures" / "figure_use_case_overview.png"
REMINDER = ROOT / "Latex" / "dissertation_figures" / "figure_reminder_workflow.png"


ACCENT = "0F766E"
ACCENT_LIGHT = "CCFBF1"
TITLE_COLOR = "0F172A"
TEXT_COLOR = "1F2937"
MUTED = "475569"


def style_title(shape, text: str) -> None:
    shape.text = text
    p = shape.text_frame.paragraphs[0]
    p.font.name = "Aptos Display"
    p.font.bold = True
    p.font.size = Pt(28)
    p.font.color.rgb = _rgb(TITLE_COLOR)


def style_body(text_frame, bullets: list[str], font_size: int = 18) -> None:
    text_frame.clear()
    for idx, bullet in enumerate(bullets):
        p = text_frame.paragraphs[0] if idx == 0 else text_frame.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.name = "Aptos"
        p.font.size = Pt(font_size)
        p.font.color.rgb = _rgb(TEXT_COLOR)
        p.space_after = Pt(8)


def _rgb(hex_value: str):
    from pptx.dml.color import RGBColor

    return RGBColor.from_string(hex_value)


def add_banner(slide, title: str, subtitle: str | None = None):
    banner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.9))
    banner.fill.solid()
    banner.fill.fore_color.rgb = _rgb(ACCENT)
    banner.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.45), Inches(0.15), Inches(8.5), Inches(0.35))
    title_p = title_box.text_frame.paragraphs[0]
    title_p.text = title
    title_p.font.name = "Aptos Display"
    title_p.font.bold = True
    title_p.font.size = Pt(26)
    title_p.font.color.rgb = _rgb("FFFFFF")

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.48), Inches(0.48), Inches(8.5), Inches(0.22))
        sub_p = sub_box.text_frame.paragraphs[0]
        sub_p.text = subtitle
        sub_p.font.name = "Aptos"
        sub_p.font.size = Pt(11)
        sub_p.font.color.rgb = _rgb("E6FFFA")


def add_footer(slide, text: str):
    box = slide.shapes.add_textbox(Inches(9.8), Inches(7.1), Inches(3), Inches(0.25))
    p = box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    p.text = text
    p.font.name = "Aptos"
    p.font.size = Pt(10)
    p.font.color.rgb = _rgb(MUTED)


def add_title_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = _rgb("F8FAFC")

    ribbon = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(0.6), Inches(12.1), Inches(0.32))
    ribbon.fill.solid()
    ribbon.fill.fore_color.rgb = _rgb(ACCENT_LIGHT)
    ribbon.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.2), Inches(10.8), Inches(1.4))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Agentic AI for Education"
    p.font.name = "Aptos Display"
    p.font.bold = True
    p.font.size = Pt(30)
    p.font.color.rgb = _rgb(TITLE_COLOR)
    p2 = tf.add_paragraph()
    p2.text = "Design, Implementation and Evaluation of a Learner-Support System"
    p2.font.name = "Aptos Display"
    p2.font.size = Pt(22)
    p2.font.color.rgb = _rgb(ACCENT)

    meta = slide.shapes.add_textbox(Inches(0.78), Inches(3.1), Inches(6.5), Inches(1.6))
    style_body(
        meta.text_frame,
        [
            "MSc Computer Science Dissertation Project",
            "Focus: active learner support through scheduling, adaptation, reminders, and explanation",
            "Core stack: React PWA, Flask, PostgreSQL, LangGraph, Supabase, Brevo",
        ],
        font_size=18,
    )

    callout = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.2), Inches(2.2), Inches(4.2), Inches(2.0))
    callout.fill.solid()
    callout.fill.fore_color.rgb = _rgb("ECFEFF")
    callout.line.color.rgb = _rgb("0891B2")
    tf2 = callout.text_frame
    p = tf2.paragraphs[0]
    p.text = "Problem"
    p.font.bold = True
    p.font.size = Pt(20)
    p.font.color.rgb = _rgb(TITLE_COLOR)
    p = tf2.add_paragraph()
    p.text = "Passive tools store tasks, but they do not actively help learners decide what to do next or recover from change."
    p.font.size = Pt(16)
    p.font.color.rgb = _rgb(TEXT_COLOR)

    add_footer(slide, "Dissertation presentation")


def add_bullets_slide(prs: Presentation, title: str, bullets: list[str], subtitle: str | None = None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _rgb("FFFFFF")
    add_banner(slide, title, subtitle)

    body = slide.shapes.add_textbox(Inches(0.75), Inches(1.35), Inches(11.8), Inches(5.6))
    style_body(body.text_frame, bullets, font_size=20)
    add_footer(slide, "Agentic AI for Education")
    return slide


def add_image_slide(prs: Presentation, title: str, image: Path, caption: str, bullets: list[str]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _rgb("FFFFFF")
    add_banner(slide, title)
    slide.shapes.add_picture(str(image), Inches(0.7), Inches(1.35), width=Inches(7.0))
    caption_box = slide.shapes.add_textbox(Inches(0.75), Inches(5.82), Inches(6.8), Inches(0.35))
    cp = caption_box.text_frame.paragraphs[0]
    cp.text = caption
    cp.font.name = "Aptos"
    cp.font.italic = True
    cp.font.size = Pt(11)
    cp.font.color.rgb = _rgb(MUTED)

    bullet_box = slide.shapes.add_textbox(Inches(8.0), Inches(1.45), Inches(4.6), Inches(4.9))
    style_body(bullet_box.text_frame, bullets, font_size=18)
    add_footer(slide, "Agentic AI for Education")


def add_results_slide(prs: Presentation, results: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _rgb("FFFFFF")
    add_banner(slide, "Evaluation Summary", "Executable scenario runs from the project evaluation framework")

    headers = ["Scenario", "Baseline", "Agent", "Interpretation"]
    rows = [
        [
            "A: Standard week",
            f"{results['scenario_a']['baseline']['avg_conflict_rate']:.1f}% conflicts / {results['scenario_a']['baseline']['avg_deadline_compliance_rate']:.1f}% compliance",
            f"{results['scenario_a']['agent']['avg_conflict_rate']:.1f}% conflicts / {results['scenario_a']['agent']['avg_deadline_compliance_rate']:.1f}% compliance",
            "Near parity under ordinary workload",
        ],
        [
            "C: Disruption",
            f"{results['scenario_c']['baseline']['avg_efficiency']:.1f}% efficiency",
            f"{results['scenario_c']['agent']['avg_efficiency']:.1f}% efficiency",
            "Current adaptation logic still converges with baseline",
        ],
        [
            "D: Dynamic addition",
            f"{results['scenario_d']['baseline']['avg_conflict_rate']:.1f}% conflicts",
            f"{results['scenario_d']['agent']['avg_conflict_rate']:.1f}% conflicts",
            "Agent improves conflict handling after urgent additions",
        ],
        [
            "E: 5-week project",
            f"{results['scenario_e']['baseline']['avg_deadline_compliance_rate']:.1f}% compliance",
            f"{results['scenario_e']['agent']['avg_deadline_compliance_rate']:.1f}% compliance",
            "Severe saturation exposes current prototype limits",
        ],
    ]

    table = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(0.6), Inches(1.45), Inches(12.0), Inches(3.3)).table
    for idx, header in enumerate(headers):
        cell = table.cell(0, idx)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = _rgb(ACCENT)
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.bold = True
            paragraph.font.size = Pt(13)
            paragraph.font.color.rgb = _rgb("FFFFFF")

    for r_idx, row in enumerate(rows, start=1):
        for c_idx, value in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = value
            cell.fill.solid()
            cell.fill.fore_color.rgb = _rgb("F8FAFC" if r_idx % 2 else "FFFFFF")
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(11)
                paragraph.font.color.rgb = _rgb(TEXT_COLOR)

    notes = slide.shapes.add_textbox(Inches(0.75), Inches(5.15), Inches(11.6), Inches(1.4))
    style_body(
        notes.text_frame,
        [
            "The evaluation evidence is strongest for dynamic-change scenarios and weakest under severe task saturation.",
            "This supports a bounded claim: the system is a working learner-support system, not a fully optimised planner.",
        ],
        font_size=18,
    )
    add_footer(slide, "Evaluation evidence")


def build_deck(results: dict) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs)
    add_bullets_slide(
        prs,
        "Problem and Research Aim",
        [
            "Learners often rely on passive tools that record deadlines but do not actively support prioritisation, scheduling, replanning, or follow-through.",
            "The project asks whether an agentic learner-support system can provide active, dependable, and testable support instead of static record-keeping.",
            "The design target is broader than one cohort: a learner-support system for education in structured study settings.",
        ],
    )
    add_bullets_slide(
        prs,
        "Core Features",
        [
            "Task capture through forms and natural language",
            "Schedule generation with deterministic fallback",
            "Automatic adaptation after task or availability changes",
            "Conflict detection and explanatory support",
            "Daily due-soon and overdue reminders with hosted automation",
            "React PWA frontend with Flask, PostgreSQL, and LangGraph backend orchestration",
        ],
    )
    add_image_slide(
        prs,
        "Use Cases",
        USE_CASE,
        "Figure: learner-facing use cases carried into implementation, testing, and evaluation",
        [
            "UC1: capture work quickly",
            "UC2: review and organise workload",
            "UC3: generate a study schedule",
            "UC4: adapt after change",
            "UC5: understand priorities and conflicts",
            "UC6: receive proactive reminder support",
        ],
    )
    add_image_slide(
        prs,
        "System Architecture",
        ARCH,
        "Figure: architecture of the learner-support system",
        [
            "React PWA handles the learner-facing interface",
            "Flask exposes deterministic APIs",
            "LangGraph routes and coordinates agent behaviour",
            "PostgreSQL stores tasks, schedules, sessions, and availability",
            "External services remain optional and bounded",
        ],
    )
    add_image_slide(
        prs,
        "Hosted Reminder Automation",
        REMINDER,
        "Figure: cron-triggered reminder flow using Supabase and Brevo",
        [
            "Supabase Cron triggers the reminder function daily",
            "The function queries due-soon and overdue pending tasks",
            "Brevo sends transactional email to the learner address on file",
            "The reminder log prevents duplicate sends on the same day/window",
        ],
    )
    add_bullets_slide(
        prs,
        "Implementation Decisions",
        [
            "LangGraph was chosen over AutoGen and CrewAI because it provided explicit state management and lower dependency on live LLM calls.",
            "The system uses a hybrid architecture: deterministic operations for CRUD and security-sensitive flows, selective AI support for planning and explanation.",
            "Supabase + Brevo provide an always-on reminder path without depending on a developer laptop.",
        ],
    )
    add_results_slide(prs, results)
    add_bullets_slide(
        prs,
        "Strengths, Limits, and Next Steps",
        [
            "Strengths: active learner support, broad workflow coverage, deployable reminder automation, and strong implementation/testing evidence.",
            "Limits: scheduling quality converges with baseline under heavy saturation; reminder delivery still depends on valid provider credentials.",
            "Next steps: richer adaptation logic, more severe scenario coverage, UI screenshots in the dissertation, and stronger end-to-end reminder verification.",
        ],
    )
    add_bullets_slide(
        prs,
        "Presentation Close",
        [
            "The project delivers a working learner-support system that goes beyond passive task storage.",
            "Its strongest contribution is the integration of capture, scheduling, adaptation, explanation, and reminders within one coherent system.",
            "The correct academic claim is a tested, implemented, and deployable learner-support system with bounded evaluation claims.",
        ],
    )

    prs.save(str(OUTPUT))

    NOTES.write_text(
        "\n".join(
            [
                "# Agentic AI for Education Presentation Notes",
                "",
                "1. Title: state the project title and the learner-support framing.",
                "2. Problem and Research Aim: stress the shift from passive tools to active learner support.",
                "3. Core Features: explain the six practical capabilities.",
                "4. Use Cases: walk quickly through UC1-UC6.",
                "5. Architecture: explain frontend, backend, LangGraph, and database roles.",
                "6. Hosted Reminder Automation: explain cron, edge function, and Brevo in simple terms.",
                "7. Implementation Decisions: justify LangGraph and the hybrid architecture.",
                "8. Evaluation Summary: highlight parity in stable cases, gains in dynamic addition, and limitations in saturation.",
                "9. Strengths, Limits, and Next Steps: keep claims bounded.",
                "10. Close: describe the work as an implemented and deployable learner-support system.",
            ]
        ),
        encoding="utf-8",
    )


def main():
    runner = TestRunner(runs_per_scenario=3)
    results = runner.run_all_scenarios()
    (ROOT / "presentation_evaluation_results.json").write_text(json.dumps(results, indent=2, default=str), encoding="utf-8")
    build_deck(results)


if __name__ == "__main__":
    main()
